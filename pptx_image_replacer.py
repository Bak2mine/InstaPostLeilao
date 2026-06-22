"""
PPTX Image Replacer - Replace template images with property photos
"""

import requests
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from io import BytesIO
import re
from pathlib import Path

class PropertyImageReplacer:
    """Replace PPTX images with property photos"""

    def __init__(self, pptx_path, output_dir):
        self.pptx_path = pptx_path
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(exist_ok=True)

    def extract_first_image_url(self, html_text):
        """Extract first image URL from property page HTML"""
        # Pattern: src="https://leiloariasmart.com.br/_admin_/upload/..."
        matches = re.findall(r'src="(https://leiloariasmart\.com\.br/_admin_/upload/[^"]+)"', html_text)
        if matches:
            return matches[0]
        return None

    def download_image(self, image_url):
        """Download image from URL"""
        try:
            response = requests.get(image_url, timeout=30)
            if response.status_code == 200:
                return Image.open(BytesIO(response.content))
        except Exception as e:
            print(f"Error downloading image: {e}")
        return None

    def replace_images_in_pptx(self, image_pil, output_pptx_path):
        """Replace images in PPTX - test by replacing the middle-sized and other candidates"""
        prs = Presentation(self.pptx_path)
        slide = prs.slides[0]

        replaced_count = 0
        test_indices = []  # Track which we replaced for testing

        for shape in slide.shapes:
            if hasattr(shape, "fill") and shape.fill.type == 6:  # PICTURE type
                width = shape.width
                height = shape.height
                width_inches = width / Inches(1)
                height_inches = height / Inches(1)

                print(f"  Found image: {shape.name} - {width_inches:.2f}\" x {height_inches:.2f}\"")

                # Strategy: Replace the wider images (likely sky/cloud and backgrounds)
                # Freeform 17, 18, 29 (4.39 x 1.22) - likely candidates
                # Freeform 13 (0.83 x 5.35) - likely tall side element
                # Try replacing Freeform 17 (and others of same size) as primary candidate

                if "Freeform 17" in shape.name or "Freeform 18" in shape.name or "Freeform 29" in shape.name:
                    print(f"    -> Testing replacement (sky/cloud candidate)")
                    self._replace_image(shape, image_pil)
                    replaced_count += 1
                    test_indices.append(shape.name)

        # Save modified PPTX
        prs.save(output_pptx_path)
        print(f"\n[OK] Replaced {replaced_count} images")
        print(f"     Replaced shapes: {', '.join(test_indices)}")
        print(f"     OPEN THE PPTX AND CHECK IF IMAGES LOOK RIGHT")
        return output_pptx_path

    def _replace_image(self, shape, image_pil):
        """Replace image in shape"""
        try:
            from pptx.util import Inches

            # Save image to bytes
            img_bytes = BytesIO()
            image_pil.save(img_bytes, format='PNG')
            img_bytes.seek(0)

            # Access image through the element
            fill = shape.fill
            if hasattr(fill, '_element'):
                # Get the blip element (image reference)
                blip_rId = fill._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}blipFill/{http://schemas.openxmlformats.org/presentationml/2006/main}blip')

                if blip_rId is not None:
                    # Get relationship ID
                    rId = blip_rId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                    if rId:
                        # Get the image part from the shape's part
                        part = shape.part
                        image_part = part.related_part(rId)
                        image_part._blob = img_bytes.getvalue()
                        print(f"       [Replaced successfully]")
                        return

            print(f"       [Could not find image reference]")

        except Exception as e:
            print(f"       [Error: {str(e)[:60]}]")

    def generate_pdf(self, pptx_path, pdf_path):
        """Convert PPTX to PDF (requires LibreOffice or similar)"""
        import subprocess
        try:
            subprocess.run([
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                "--headless",
                "--convert-to", "pdf",
                "--outdir", str(self.output_dir),
                str(pptx_path)
            ], timeout=60, check=True)
            print(f"✓ Generated PDF: {pdf_path}")
            return True
        except Exception as e:
            print(f"Warning: Could not convert to PDF: {e}")
            print(f"  PPTX saved at: {pptx_path}")
            return False


def main():
    print("=" * 60)
    print("PPTX Image Replacer - Test")
    print("=" * 60)

    # Test with Casa 147m² property
    test_url = "https://leiloariasmart.com.br/imovel/1549"
    pptx_path = r"C:\Users\andre\Desktop\Leiloaria\Post\1 e 2 praça.pptx"
    output_dir = r"C:\Users\andre\Desktop\Leiloaria\Post\test_output"

    replacer = PropertyImageReplacer(pptx_path, output_dir)

    print(f"\n[1/3] Downloading property page...")
    headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'}
    response = requests.get(test_url, headers=headers, timeout=30)
    response.encoding = 'utf-8'
    html = response.text

    print(f"[2/3] Extracting first image...")
    image_url = replacer.extract_first_image_url(html)
    if not image_url:
        print("ERROR: No image found!")
        return

    print(f"    Image URL: {image_url}")
    image = replacer.download_image(image_url)
    if not image:
        print("ERROR: Could not download image!")
        return

    print(f"    Image size: {image.size}")

    print(f"[3/3] Replacing images in PPTX...")
    output_pptx = Path(output_dir) / "PPTX_WITH_IMAGES.pptx"
    replacer.replace_images_in_pptx(image, str(output_pptx))

    print("\n" + "=" * 60)
    print("Complete!")
    print(f"  PPTX: {output_pptx}")
    print("  Open in PowerPoint to verify images were replaced")


if __name__ == "__main__":
    main()
