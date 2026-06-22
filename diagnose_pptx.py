"""
Diagnostic tool to inspect PPTX structure and identify shapes/groups
"""

from pptx import Presentation
from pptx.util import Inches
import sys

def diagnose_pptx(pptx_path):
    """Print detailed information about PPTX structure"""

    prs = Presentation(pptx_path)

    print("=" * 80)
    print(f"PPTX File: {pptx_path}")
    print("=" * 80)

    print(f"\nSlide dimensions: {prs.slide_width / 914400:.2f} x {prs.slide_height / 914400:.2f} inches")
    print(f"Total slides: {len(prs.slides)}")

    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"\n{'='*80}")
        print(f"SLIDE {slide_num}")
        print(f"{'='*80}")
        print(f"Total shapes: {len(slide.shapes)}\n")

        for shape_num, shape in enumerate(slide.shapes, 1):
            try:
                shape_name = shape.name
                shape_type = str(shape.shape_type)

                # Get dimensions
                width_in = shape.width / 914400
                height_in = shape.height / 914400
                left_in = shape.left / 914400
                top_in = shape.top / 914400

                # Check if it has text
                has_text = False
                text_content = ""
                if hasattr(shape, 'text'):
                    text_content = shape.text[:50] if shape.text else ""
                    has_text = bool(shape.text.strip())

                # Check if it's a group or picture
                is_group = "Group" in shape_type
                is_picture = "Picture" in shape_type

                print(f"Shape {shape_num}:")
                print(f"  Name: {shape_name}")
                print(f"  Type: {shape_type}")
                print(f"  Position: ({left_in:.2f}, {top_in:.2f}) inches")
                print(f"  Size: {width_in:.2f} x {height_in:.2f} inches")

                if has_text:
                    print(f"  Text: {text_content}")

                # If it's a group, show what's inside
                if is_group:
                    print(f"  GROUP CONTENTS:")
                    try:
                        for sub_num, sub_shape in enumerate(shape.shapes, 1):
                            sub_name = sub_shape.name
                            sub_type = str(sub_shape.shape_type)
                            sub_width = sub_shape.width / 914400
                            sub_height = sub_shape.height / 914400
                            print(f"    {sub_num}. {sub_name} ({sub_type}) - {sub_width:.2f}x{sub_height:.2f}in")
                    except:
                        print(f"    (Could not enumerate group contents)")

                if is_picture:
                    print(f"  PICTURE")

                print()

            except Exception as e:
                print(f"Shape {shape_num}: Error inspecting - {e}\n")

if __name__ == "__main__":
    # Diagnose all templates
    templates = [
        r"C:\Users\andre\Desktop\Leiloaria\Post\1 e 2 praça.pptx",
        r"C:\Users\andre\Desktop\Leiloaria\Post\2praca.pptx",
        r"C:\Users\andre\Desktop\Leiloaria\Post\1prac.pptx",
        r"C:\Users\andre\Desktop\Leiloaria\Post\1pracD.pptx",
    ]

    for template in templates:
        try:
            print(f"\n\n")
            diagnose_pptx(template)
        except Exception as e:
            print(f"Error with {template}: {e}")

    print("\n\n" + "=" * 80)
    print("DIAGNOSTIC COMPLETE")
    print("=" * 80)
