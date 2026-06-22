"""
Final PPTX Generator - Text + Image replacement with PDF export
Complete automation: scrape → modify text → replace images → export PDF
"""

import requests
from pptx import Presentation
from PIL import Image
from io import BytesIO
import re
from pathlib import Path
import zipfile
import shutil
import subprocess

class PropertyPDFGenerator:
    """Complete property PDF generator"""

    def __init__(self, template_pptx, link_pdf, output_dir):
        self.template_pptx = template_pptx
        self.link_pdf = link_pdf
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(exist_ok=True)

    def extract_first_image_url(self, html_text):
        """Extract first image URL from property page HTML"""
        matches = re.findall(r'src="(https://leiloariasmart\.com\.br/_admin_/upload/[^"]+)"', html_text)
        return matches[0] if matches else None

    def download_image(self, image_url):
        """Download image from URL"""
        try:
            response = requests.get(image_url, timeout=30)
            if response.status_code == 200:
                return Image.open(BytesIO(response.content))
        except Exception as e:
            print(f"Error downloading image: {e}")
        return None

    def modify_pptx_text(self, property_data, output_pptx):
        """Modify text in PPTX while preserving formatting"""
        prs = Presentation(self.template_pptx)
        slide = prs.slides[0]

        # Mapping of old text to new text
        praca1_data = property_data.get('praca1_data', '')
        praca2_data = property_data.get('praca2_data', '')
        praca1_date_short = praca1_data.split('/')[0] + '/' + praca1_data.split('/')[1] if praca1_data else '15/06'
        praca2_date_short = praca2_data.split('/')[0] + '/' + praca2_data.split('/')[1] if praca2_data else '30/06'

        desconto_pct = property_data.get('desconto_pct')
        desconto_text = f"{int(desconto_pct)}% DE DESCONTO!" if desconto_pct else "20% DE DESCONTO!"

        # Handle both template styles (1prac uses different placeholder)
        new_title = property_data.get('titulo', '')
        replacements = {
            "Lote de Terreno 300m²": new_title,  # For 2praca and "1 e 2 praça" templates
            "Terrenos em Cond. até 566m²": new_title,  # For 1prac and 1pracD templates
            "NOVA ODESSA/SP": f"{property_data.get('cidade', '')}/{property_data.get('estado', '')}",
            "TERESINA/PI": f"{property_data.get('cidade', '')}/{property_data.get('estado', '')}",  # Alternative location placeholder
            "R$ 255.056,73": property_data.get('praca1_valor', '') or 'R$ 0,00',
            "R$ 142.920,39": property_data.get('praca1_valor', '') or 'R$ 0,00',  # For 1prac templates
            "R$ 245.848,12": property_data.get('praca2_valor', '') or 'R$ 0,00',
            "15/06": praca1_date_short,
            "30/06": praca2_date_short,
            "02/07": praca1_date_short,  # For 1prac templates
            "20% DE DESCONTO!": desconto_text,
        }

        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                original_text = shape.text.strip()

                if original_text in replacements:
                    new_text = replacements[original_text]
                    text_frame = shape.text_frame

                    # Disable word wrap for all fields to prevent text hanging
                    text_frame.word_wrap = False

                    for paragraph in text_frame.paragraphs:
                        if paragraph.runs:
                            first_run = paragraph.runs[0]
                            font_copy = first_run.font

                            for run in paragraph.runs:
                                run.text = ""

                            new_run = paragraph.add_run()
                            new_run.text = new_text

                            # Preserve formatting
                            new_run.font.name = font_copy.name
                            new_run.font.size = font_copy.size
                            new_run.font.bold = font_copy.bold
                            new_run.font.italic = font_copy.italic
                            if font_copy.color.rgb:
                                new_run.font.color.rgb = font_copy.color.rgb

        prs.save(output_pptx)
        return output_pptx

    def replace_images_in_pptx(self, image_pil, pptx_path):
        """Replace sky placeholder + add scaled background (crop sidebars)"""
        try:
            # Crop image to remove cyan/blue sidebars (crop 20% from left and right)
            width, height = image_pil.size
            crop_left = int(width * 0.20)
            crop_right = int(width * 0.80)
            image_cropped = image_pil.crop((crop_left, 0, crop_right, height))

            temp_dir = self.output_dir / "temp_pptx"
            if temp_dir.exists():
                shutil.rmtree(temp_dir)
            temp_dir.mkdir()

            with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
                zip_ref.extractall(temp_dir)

            image_rgb = image_cropped.convert('RGB')

            # Replace Freeform 6 (sky image) - image1.jpeg
            sky_image_path = temp_dir / "ppt" / "media" / "image1.jpeg"
            if sky_image_path.exists():
                image_rgb.save(str(sky_image_path))

            # Re-zip the PPTX
            with zipfile.ZipFile(pptx_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
                for root, dirs, files in __import__('os').walk(temp_dir):
                    for file in files:
                        file_path = Path(root) / file
                        arcname = file_path.relative_to(temp_dir)
                        zipf.write(file_path, arcname)

            # Now add scaled background image
            prs = Presentation(pptx_path)
            slide = prs.slides[0]

            from pptx.util import Inches

            # Delete Group 2 (gray background with blue/cyan borders)
            for shape in list(slide.shapes):
                if shape.name == "Group 2":
                    sp = shape.element
                    sp.getparent().remove(sp)
                    break

            # Scale background to full slide dimensions
            scaled_width = Inches(11.25)
            scaled_height = Inches(14.06)

            # Position at top-left (no offset)
            left = Inches(0)
            top = Inches(0)

            img_stream = BytesIO()
            image_rgb.save(img_stream, format='JPEG')
            img_stream.seek(0)

            picture_bg = slide.shapes.add_picture(img_stream, left, top, width=scaled_width, height=scaled_height)
            slide.shapes._spTree.remove(picture_bg._element)
            slide.shapes._spTree.insert(2, picture_bg._element)

            prs.save(str(pptx_path))
            print("    Images replaced (sky) + added (background, scaled)")
            return True

        except Exception as e:
            print(f"    Error replacing images: {e}")
            return False
        finally:
            if temp_dir.exists():
                shutil.rmtree(temp_dir)

    def export_to_pdf(self, pptx_path, pdf_path):
        """Export PPTX to PDF using LibreOffice"""
        try:
            subprocess.run([
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                "--headless",
                "--convert-to", "pdf",
                "--outdir", str(self.output_dir),
                str(pptx_path)
            ], timeout=60, check=True)

            # Rename the output PDF
            generated_pdf = self.output_dir / f"{pptx_path.stem}.pdf"
            if generated_pdf.exists():
                generated_pdf.rename(pdf_path)
            print(f"    Exported to PDF")
            return True

        except Exception as e:
            print(f"    Warning: PDF export failed: {e}")
            return False

    def generate(self, property_url, property_data, output_name):
        """Complete generation pipeline"""
        print(f"\n[1/4] Downloading property page...")
        headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'}

        try:
            response = requests.get(property_url, headers=headers, timeout=30)
            response.encoding = 'utf-8'
            html = response.text
        except Exception as e:
            print(f"    Error: {e}")
            return None

        print(f"[2/4] Extracting first image...")
        image_url = self.extract_first_image_url(html)
        if not image_url:
            print("    Warning: No image found on property page")
            image = None
        else:
            image = self.download_image(image_url)
            if image:
                print(f"    Downloaded: {image.size}")

        print(f"[3/4] Modifying PPTX text...")
        pptx_path = self.output_dir / f"{output_name}.pptx"
        self.modify_pptx_text(property_data, str(pptx_path))

        print(f"[4/4] Replacing images...")
        if image:
            self.replace_images_in_pptx(image, pptx_path)
        else:
            print("    Skipped (no image downloaded)")

        print(f"[5/5] Adding page 2 with background image...")
        if image:
            self.add_link_page_with_image(image, pptx_path)

        print(f"    Done!")
        return pptx_path

    def add_link_page_with_image(self, image_pil, pptx_path):
        """Add background image to slide 2 (LINK page)"""
        try:
            prs = Presentation(str(pptx_path))

            # Use existing slide 2 (LINK page)
            if len(prs.slides) < 2:
                return True  # No slide 2 to modify

            slide = prs.slides[1]

            from pptx.util import Inches

            # Crop image to remove sidebars (20% from each side)
            width, height = image_pil.size
            crop_left = int(width * 0.20)
            crop_right = int(width * 0.80)
            image_cropped = image_pil.crop((crop_left, 0, crop_right, height))

            # Add background image to slide 2
            image_rgb = image_cropped.convert('RGB')
            img_stream = BytesIO()
            image_rgb.save(img_stream, format='JPEG')
            img_stream.seek(0)

            picture_bg = slide.shapes.add_picture(img_stream, 0, 0, width=Inches(11.25), height=Inches(14.06))
            slide.shapes._spTree.remove(picture_bg._element)
            slide.shapes._spTree.insert(2, picture_bg._element)

            # Add "LINK" text in the center
            left = Inches(3.5)
            top = Inches(6.0)
            width = Inches(4.25)
            height = Inches(2.0)

            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = "LINK"

            # Format the text
            from pptx.util import Pt
            from pptx.dml.color import RGBColor

            for paragraph in text_frame.paragraphs:
                paragraph.alignment = 1  # Center align
                for run in paragraph.runs:
                    run.font.size = Pt(60)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(255, 255, 255)  # White

            prs.save(str(pptx_path))
            return True

        except Exception as e:
            print(f"    Error adding link page: {e}")
            return False


def main():
    print("=" * 60)
    print("Property PDF Generator - Test")
    print("=" * 60)

    # Test with Casa 147m² property
    generator = PropertyPDFGenerator(
        template_pptx=r"C:\Users\andre\Desktop\Leiloaria\Post\1 e 2 praça.pptx",
        link_pdf=r"C:\Users\andre\Desktop\Leiloaria\Post\Posts IG.pdf",
        output_dir=r"C:\Users\andre\Desktop\Leiloaria\Post\test_output"
    )

    property_data = {
        'titulo': 'Casa 147m²',
        'cidade': 'Uberlandia',
        'estado': 'MG',
        'praca1_valor': 'R$ 320.418,61',
        'praca2_valor': 'R$ 860.288,76',
        'praca1_data': '15/06/2026',
        'praca2_data': '29/06/2026',
        'desconto_pct': 36,
    }

    pdf = generator.generate(
        property_url="https://leiloariasmart.com.br/imovel/1549",
        property_data=property_data,
        output_name="TEST_Casa_147m"
    )

    print("\n" + "=" * 60)
    if pdf and pdf.exists():
        print(f"Success! PPTX created with text + images:")
        print(f"  {pdf}")
        print(f"\nTo generate PDF: Open in PowerPoint and Export as PDF")
    else:
        print("Generation had issues - check output above")


if __name__ == "__main__":
    main()
