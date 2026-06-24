"""
PowerPoint template handling - text and image replacement
"""

import zipfile
import shutil
import os
from pathlib import Path
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
from typing import Optional, Dict
from config import SLIDE_WIDTH_INCHES, SLIDE_HEIGHT_INCHES

class PPTXHandler:
    """Handle PPTX template modifications - text and image replacement"""

    @staticmethod
    def replace_embedded_images(pptx_path: Path, image_pil: Image.Image) -> bool:
        """NOT USED - Embedded images are left alone to avoid affecting logos/decorations.
        The middle image will come through the background via the actual slide image addition.

        Args:
            pptx_path: Path to PPTX file
            image_pil: PIL Image to use as replacement

        Returns:
            True (no-op, just skipped)
        """
        print(f"    Skipping embedded image replacement (using background image instead)")
        return True

    @staticmethod
    def truncate_title_for_overflow(title: str, max_title_length: int = 50) -> str:
        """Truncate title to prevent overflow in PowerPoint text boxes.

        Args:
            title: Title text (should not contain city/state)
            max_title_length: Maximum characters allowed for title

        Returns:
            Truncated title if needed, otherwise original title
        """
        if len(title) > max_title_length:
            # Cut at max length and remove trailing incomplete words
            truncated = title[:max_title_length].rstrip()
            last_space = truncated.rfind(' ')
            if last_space > 0:
                truncated = truncated[:last_space]
            # Preserve m² unit if present
            if 'm²' in title:
                truncated = truncated.rstrip('m').rstrip() + ' m²'
            return truncated.strip()
        return title

    @staticmethod
    def modify_text(template_pptx: Path, property_data: Dict, output_pptx: Path) -> Path:
        """Modify text in PPTX while preserving formatting

        Args:
            template_pptx: Path to template PPTX file
            property_data: Property data dict with keys like 'titulo', 'cidade', etc
            output_pptx: Path to save modified PPTX

        Returns:
            Path to output PPTX
        """
        prs = Presentation(str(template_pptx))
        slide = prs.slides[0]

        # Calculate date shorthand (MM/DD from MM/DD/YYYY)
        praca1_data = property_data.get('praca1_data', '')
        praca2_data = property_data.get('praca2_data', '')
        praca1_date_short = praca1_data.split('/')[0] + '/' + praca1_data.split('/')[1] if praca1_data else '15/06'
        praca2_date_short = praca2_data.split('/')[0] + '/' + praca2_data.split('/')[1] if praca2_data else '30/06'

        desconto_pct = property_data.get('desconto_pct')
        desconto_text = f"{int(desconto_pct)}% DE DESCONTO!" if desconto_pct else "20% DE DESCONTO!"

        # Text replacement mapping for all template types
        raw_title = property_data.get('titulo', '')
        location = f"{property_data.get('cidade', '')}/{property_data.get('estado', '')}"

        # Truncate title if it's too long for the PowerPoint text box
        new_title = PPTXHandler.truncate_title_for_overflow(raw_title)

        replacements = {
            # Title placeholders (different templates have different defaults)
            "Lote de Terreno 300m²": new_title,
            "Terrenos em Cond. até 566m²": new_title,
            "Casa 147m²": new_title,
            "Apartamento 102m²": new_title,
            # Location placeholders
            "NOVA ODESSA/SP": location,
            "TERESINA/PI": location,
            "UBERLÂNDIA/MG": location,
            # Price placeholders (both auctions)
            "R$ 255.056,73": property_data.get('praca1_valor', '') or 'R$ 0,00',
            "R$ 142.920,39": property_data.get('praca1_valor', '') or 'R$ 0,00',
            "R$ 245.848,12": property_data.get('praca2_valor', '') or 'R$ 0,00',
            # Date placeholders
            "15/06": praca1_date_short,
            "02/07": praca1_date_short,
            "30/06": praca2_date_short,
            # Discount
            "20% DE DESCONTO!": desconto_text,
        }

        # Replace text in all shapes on slide 1
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                original_text = shape.text.strip()

                if original_text in replacements:
                    new_text = replacements[original_text]
                    text_frame = shape.text_frame

                    # Disable word wrap to prevent text hanging
                    text_frame.word_wrap = False

                    for paragraph in text_frame.paragraphs:
                        if paragraph.runs:
                            first_run = paragraph.runs[0]
                            font_copy = first_run.font

                            # Clear existing runs
                            for run in paragraph.runs:
                                run.text = ""

                            # Add new text with preserved formatting
                            new_run = paragraph.add_run()
                            new_run.text = new_text

                            # Preserve formatting from original
                            new_run.font.name = font_copy.name
                            new_run.font.size = font_copy.size
                            new_run.font.bold = font_copy.bold
                            new_run.font.italic = font_copy.italic
                            if font_copy.color.rgb:
                                new_run.font.color.rgb = font_copy.color.rgb

        prs.save(str(output_pptx))
        return output_pptx

    @staticmethod
    def replace_images(pptx_path: Path, image_pil: Image.Image) -> bool:
        """Replace placeholder images via ZIP and delete overlaying elements.

        Args:
            pptx_path: Path to PPTX file to modify
            image_pil: PIL Image object with property image

        Returns:
            True if successful, False otherwise
        """
        try:
            import tempfile
            from image_processor import ImageProcessor

            # Process image (crop + RGB conversion)
            image_cropped = ImageProcessor.crop_sidebars(image_pil)
            image_rgb = ImageProcessor.to_rgb(image_cropped)

            # Resize to match template dimensions (1438x560)
            image_rgb = image_rgb.resize((1438, 560), Image.Resampling.LANCZOS)

            print(f"    Image size after processing: {image_rgb.size}")

            # Step 1: Replace embedded images via ZIP
            with tempfile.TemporaryDirectory() as temp_dir:
                temp_dir = Path(temp_dir)

                # Extract PPTX
                with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
                    zip_ref.extractall(temp_dir)

                media_dir = temp_dir / "ppt" / "media"

                # Check which image files exist to determine template type
                has_image1_jpeg = (media_dir / "image1.jpeg").exists()
                is_two_auctions = has_image1_jpeg

                print(f"    Template type: {'2 praças' if is_two_auctions else '1 praça'}")

                # Replace the appropriate placeholder image
                if is_two_auctions:
                    # For 2 praças: replace image1.jpeg
                    image1_path = media_dir / "image1.jpeg"
                    if image1_path.exists():
                        image_rgb.save(str(image1_path), quality=95)
                        print(f"    Replaced image1.jpeg")
                else:
                    # For 1 praça: replace image2.jpeg
                    image2_path = media_dir / "image2.jpeg"
                    if image2_path.exists():
                        image_rgb.save(str(image2_path), quality=95)
                        print(f"    Replaced image2.jpeg")

                # Re-zip the PPTX
                with zipfile.ZipFile(pptx_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
                    for root, dirs, files in os.walk(temp_dir):
                        for file in files:
                            file_path = Path(root) / file
                            arcname = file_path.relative_to(temp_dir)
                            zipf.write(file_path, arcname)

            # Open PPTX and handle slides
            prs = Presentation(str(pptx_path))

            print(f"    Processing slide 1...")
            slide1 = prs.slides[0]

            # Delete Group 2 (the gray background overlay on slide 1)
            for shape in list(slide1.shapes):
                if shape.name == "Group 2":
                    try:
                        sp = shape.element
                        sp.getparent().remove(sp)
                        print(f"    Deleted Group 2 from slide 1")
                        break
                    except Exception as e:
                        pass

            # Add full-page background image to slide 1
            img_stream1 = BytesIO()
            image_rgb.save(img_stream1, format='JPEG', quality=95)
            img_stream1.seek(0)

            # For 2-praça: full page background. For 1-praça: smaller background (9" width)
            if is_two_auctions:
                bg_width = Inches(11.25)
                bg_height = Inches(14.06)
                bg_left = Inches(0)
                bg_top = Inches(0)
            else:
                # For 1-praça: center a 9" wide image to not cover the embedded image
                bg_width = Inches(9.0)
                aspect_ratio = image_rgb.width / image_rgb.height
                bg_height = Inches(9.0 / aspect_ratio)
                bg_left = Inches((11.25 - 9.0) / 2)
                bg_top = Inches((14.06 - (9.0 / aspect_ratio)) / 2)

            picture_bg1 = slide1.shapes.add_picture(
                img_stream1,
                bg_left,
                bg_top,
                width=bg_width,
                height=bg_height
            )

            # Move to position 2 (behind other elements but in front of base)
            slide1.shapes._spTree.remove(picture_bg1._element)
            slide1.shapes._spTree.insert(2, picture_bg1._element)
            print(f"    Added background image to slide 1")

            # Process Slide 2 (if it exists)
            if len(prs.slides) >= 2:
                print(f"    Processing slide 2...")
                slide2 = prs.slides[1]

                # Delete Group 2 (gray background overlay)
                for shape in list(slide2.shapes):
                    if shape.name == "Group 2":
                        try:
                            sp = shape.element
                            sp.getparent().remove(sp)
                            print(f"    Deleted Group 2 from slide 2")
                            break
                        except Exception as e:
                            pass

                # Delete all Freeforms on slide 2 (decorative lines and background images)
                # Keep Groups (they're the price/discount frames) and TextBoxes
                deleted_count = 0
                for shape in list(slide2.shapes):
                    if 'Freeform' in shape.name:
                        try:
                            sp = shape.element
                            sp.getparent().remove(sp)
                            deleted_count += 1
                        except Exception as e:
                            pass
                if deleted_count > 0:
                    print(f"    Deleted {deleted_count} freeforms from slide 2")

                # Add background image to slide 2
                img_stream2 = BytesIO()
                image_rgb.save(img_stream2, format='JPEG', quality=95)
                img_stream2.seek(0)

                # For 2-praça: full page. For 1-praça: smaller (9" width)
                if is_two_auctions:
                    bg_width_s2 = Inches(11.25)
                    bg_height_s2 = Inches(14.06)
                    bg_left_s2 = Inches(0)
                    bg_top_s2 = Inches(0)
                else:
                    bg_width_s2 = Inches(9.0)
                    aspect_ratio = image_rgb.width / image_rgb.height
                    bg_height_s2 = Inches(9.0 / aspect_ratio)
                    bg_left_s2 = Inches((11.25 - 9.0) / 2)
                    bg_top_s2 = Inches((14.06 - (9.0 / aspect_ratio)) / 2)

                picture_bg2 = slide2.shapes.add_picture(
                    img_stream2,
                    bg_left_s2,
                    bg_top_s2,
                    width=bg_width_s2,
                    height=bg_height_s2
                )

                # Move to position 2 (behind frames/text)
                slide2.shapes._spTree.remove(picture_bg2._element)
                slide2.shapes._spTree.insert(2, picture_bg2._element)
                print(f"    Added background image to slide 2")

            # Save presentation
            prs.save(str(pptx_path))
            print(f"    All slides complete, saved")

            return True

        except Exception as e:
            print(f"[ERROR] Error replacing images: {e}")
            import traceback
            traceback.print_exc()
            return False

    @staticmethod
    def add_link_page(pptx_path: Path, image_pil: Image.Image) -> bool:
        """NO-OP: Both slides are now handled in replace_images() in a single pass.

        Args:
            pptx_path: Path to PPTX file
            image_pil: PIL Image object

        Returns:
            True (skipped, already done)
        """
        print(f"    Skipping add_link_page (already handled in replace_images)")
        return True
