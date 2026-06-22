"""
Test: Replace images in the correct groups (Group 2 and Group 5)
"""

import requests
from pptx import Presentation
from PIL import Image
from io import BytesIO
import re
import zipfile
import shutil
from pathlib import Path

# Download test image
url = "https://leiloariasmart.com.br/imovel/1549"
headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'}
response = requests.get(url, headers=headers)
html = response.text

# Extract first image URL
matches = re.findall(r'src="(https://leiloariasmart\.com\.br/_admin_/upload/[^"]+)"', html)
image_url = matches[0]

print(f"Downloading: {image_url}")
img_response = requests.get(image_url)
image_pil = Image.open(BytesIO(img_response.content))
print(f"Image size: {image_pil.size}")

# Open template
template_pptx = r"C:\Users\andre\Desktop\Leiloaria\Post\1 e 2 praça.pptx"
test_output = r"C:\Users\andre\Desktop\Leiloaria\Post\test_output\TEST_IMAGE_FIX.pptx"

prs = Presentation(template_pptx)
slide = prs.slides[0]

print("\nFinding Group 2 and Group 5...")

for idx, shape in enumerate(slide.shapes):
    if hasattr(shape, 'shapes'):  # Is a group
        width = shape.width / 914400
        height = shape.height / 914400

        if "Group 2" in shape.name and 12 < width < 13 and 14 < height < 15:
            print(f"\nFound: {shape.name} ({width:.1f}\" x {height:.1f}\")")
            print(f"  This is the GRAY BACKGROUND")
            for sub_idx, sub_shape in enumerate(shape.shapes):
                if hasattr(sub_shape, "fill") and sub_shape.fill.type == 6:
                    print(f"  Contains PICTURE: {sub_shape.name}")

        if "Group 5" in shape.name and 8.5 < width < 9.5 and 5.5 < height < 6.5:
            print(f"\nFound: {shape.name} ({width:.1f}\" x {height:.1f}\")")
            print(f"  This is the SKY/CLOUD")
            for sub_idx, sub_shape in enumerate(shape.shapes):
                if hasattr(sub_shape, "fill") and sub_shape.fill.type == 6:
                    print(f"  Contains PICTURE: {sub_shape.name}")

print("\nTest complete - verify the groups above are correct")
print("Then we'll fix the replacement logic to target these specific groups")
